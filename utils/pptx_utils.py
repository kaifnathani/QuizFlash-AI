from io import BytesIO
from typing import List
from pptx import Presentation


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a .pptx file bytes using python-pptx."""
    lines = []
    try:
        prs = Presentation(BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        lines.append(text)
    except Exception:
        return ""

    return "\n".join(lines)
